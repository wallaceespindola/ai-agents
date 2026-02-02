"""
PowerPoint Generator - Create .pptx files from article content

Generates professional PowerPoint presentations from article markdown
using python-pptx library with configurable themes and templates.
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass


@dataclass
class ColorScheme:
    """Color configuration for presentation"""
    background: str = "#FFFFFF"
    text: str = "#333333"
    primary: str = "#2E5090"
    secondary: str = "#FF6B6B"
    accent: str = "#4ECDC4"
    code_background: str = "#F5F5F5"


@dataclass
class PptxConfig:
    """Configuration for PowerPoint generation"""
    aspect_ratio: str = "16:9"
    dpi: int = 300
    default_font: str = "San Francisco"
    code_font: str = "Monaco"
    theme: str = "technical"
    include_footer: bool = True
    include_slide_numbers: bool = True
    page_width: float = 10.0  # inches
    page_height: float = 5.625  # inches for 16:9
    margin: float = 0.5  # inches


class PowerPointGenerator:
    """Generate PowerPoint presentations from slide content"""

    def __init__(self, config: Optional[PptxConfig] = None):
        self.config = config or PptxConfig()
        self.color_schemes = {
            'light': ColorScheme(
                background="#FFFFFF",
                text="#333333",
                primary="#2E5090",
                secondary="#FF6B6B",
                accent="#4ECDC4",
                code_background="#F5F5F5"
            ),
            'dark': ColorScheme(
                background="#1E1E1E",
                text="#FFFFFF",
                primary="#4ECDC4",
                secondary="#FF6B6B",
                accent="#FFE66D",
                code_background="#2D2D2D"
            ),
            'technical': ColorScheme(
                background="#0D1117",
                text="#E6EDF3",
                primary="#58A6FF",
                secondary="#79C0FF",
                accent="#79C0FF",
                code_background="#0D1117"
            )
        }

    def generate_presentation(
        self,
        slides_data: List[Dict[str, Any]],
        output_path: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> bool:
        """
        Generate PowerPoint presentation from slide data

        Args:
            slides_data: List of slide dictionaries from slide generator
            output_path: Path where to save the .pptx file
            metadata: Optional presentation metadata (title, author, etc.)

        Returns:
            True if successful, False otherwise
        """

        try:
            # Try to import python-pptx
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor

        except ImportError:
            print("ERROR: python-pptx not installed")
            print("Install with: pip install python-pptx")
            return False

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(self.config.page_width)
        prs.slide_height = Inches(self.config.page_height)

        # Set metadata
        if metadata:
            prs.core_properties.title = metadata.get('title', 'Presentation')
            prs.core_properties.author = metadata.get('author', '')
            prs.core_properties.subject = metadata.get('subject', '')

        # Get color scheme
        colors = self.color_schemes.get(
            self.config.theme,
            self.color_schemes['technical']
        )

        # Add slides
        for i, slide_data in enumerate(slides_data):
            self._add_slide(prs, slide_data, colors, i + 1, len(slides_data))

        # Create output directory if needed
        output_dir = Path(output_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)

        # Save presentation
        try:
            prs.save(output_path)
            print(f"✅ Created PowerPoint: {output_path}")
            return True
        except Exception as e:
            print(f"❌ Error saving PowerPoint: {e}")
            return False

    def _add_slide(
        self,
        prs,
        slide_data: Dict[str, Any],
        colors: ColorScheme,
        slide_num: int,
        total_slides: int
    ) -> None:
        """Add a single slide to presentation"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor

        slide_type = slide_data.get('type', 'content')
        content = slide_data.get('content', {})

        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(colors.background)

        # Add content based on slide type
        if slide_type == 'title':
            self._add_title_slide(slide, content, colors)
        elif slide_type == 'content':
            self._add_content_slide(slide, content, colors)
        elif slide_type == 'code':
            self._add_code_slide(slide, content, colors)
        elif slide_type == 'visual':
            self._add_visual_slide(slide, content, colors)
        elif slide_type == 'conclusion':
            self._add_conclusion_slide(slide, content, colors)

        # Add footer
        if self.config.include_footer:
            self._add_footer(slide, colors, slide_num, total_slides)

    def _add_title_slide(
        self,
        slide,
        content: Dict[str, Any],
        colors: ColorScheme
    ) -> None:
        """Add title slide"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = content.get('title', 'Untitled')
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = self._hex_to_rgb(colors.primary)
        title_p.alignment = PP_ALIGN.CENTER

        # Author
        author_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2), Inches(9), Inches(0.5)
        )
        author_frame = author_box.text_frame
        author_p = author_frame.paragraphs[0]
        author_p.text = content.get('author', 'Author')
        author_p.font.size = Pt(24)
        author_p.font.color.rgb = self._hex_to_rgb(colors.text)
        author_p.alignment = PP_ALIGN.CENTER

        # Date and tags
        meta_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(9), Inches(1)
        )
        meta_frame = meta_box.text_frame
        meta_p = meta_frame.paragraphs[0]
        date_str = content.get('date', '')
        tags_str = ', '.join(content.get('tags', []))
        meta_p.text = f"{date_str} | {tags_str}"
        meta_p.font.size = Pt(12)
        meta_p.font.color.rgb = self._hex_to_rgb(colors.secondary)
        meta_p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        slide,
        content: Dict[str, Any],
        colors: ColorScheme
    ) -> None:
        """Add content/body slide"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor

        # Heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        heading_frame = heading_box.text_frame
        heading_p = heading_frame.paragraphs[0]
        heading_p.text = content.get('heading', 'Section')
        heading_p.font.size = Pt(40)
        heading_p.font.bold = True
        heading_p.font.color.rgb = self._hex_to_rgb(colors.primary)

        # Body text
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(9), Inches(2)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_p = body_frame.paragraphs[0]
        body_p.text = content.get('body', '')
        body_p.font.size = Pt(18)
        body_p.font.color.rgb = self._hex_to_rgb(colors.text)

        # Bullet points
        bullet_points = content.get('bullet_points', [])
        if bullet_points:
            bullet_y = 3.6
            for bullet in bullet_points[:5]:  # Max 5 bullets
                bullet_box = slide.shapes.add_textbox(
                    Inches(1), Inches(bullet_y), Inches(8.5), Inches(0.3)
                )
                bullet_frame = bullet_box.text_frame
                bullet_p = bullet_frame.paragraphs[0]
                bullet_p.text = f"• {bullet}"
                bullet_p.font.size = Pt(14)
                bullet_p.font.color.rgb = self._hex_to_rgb(colors.text)
                bullet_y += 0.35

    def _add_code_slide(
        self,
        slide,
        content: Dict[str, Any],
        colors: ColorScheme
    ) -> None:
        """Add code example slide"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor

        # Language label
        lang_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(3), Inches(0.3)
        )
        lang_frame = lang_box.text_frame
        lang_p = lang_frame.paragraphs[0]
        lang_p.text = f"Code: {content.get('language', 'text').upper()}"
        lang_p.font.size = Pt(10)
        lang_p.font.color.rgb = self._hex_to_rgb(colors.secondary)

        # Code block background
        code_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1), Inches(9), Inches(4)
        )
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = self._hex_to_rgb(colors.code_background)
        code_shape.line.color.rgb = self._hex_to_rgb(colors.accent)

        # Code text
        code_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), Inches(8.6), Inches(3.6)
        )
        code_frame = code_box.text_frame
        code_frame.word_wrap = True
        code_frame.vertical_anchor = MSO_ANCHOR.TOP

        code_p = code_frame.paragraphs[0]
        code_text = content.get('code', '')
        # Truncate if too long
        lines = code_text.split('\n')
        if len(lines) > 20:
            code_text = '\n'.join(lines[:20]) + '\n...'

        code_p.text = code_text
        code_p.font.name = self.config.code_font
        code_p.font.size = Pt(10)
        code_p.font.color.rgb = self._hex_to_rgb(colors.text)
        code_p.font.monospace = True

    def _add_visual_slide(
        self,
        slide,
        content: Dict[str, Any],
        colors: ColorScheme
    ) -> None:
        """Add visual/image slide"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # Try to add image
        file_path = content.get('file_path', '')
        if file_path and Path(file_path).exists():
            try:
                slide.shapes.add_picture(
                    file_path,
                    Inches(0.5), Inches(0.5), width=Inches(9)
                )
            except Exception:
                # If image fails, just show alt text
                pass

        # Caption
        alt_text = content.get('alt_text', '')
        if alt_text:
            caption_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.8), Inches(9), Inches(0.7)
            )
            caption_frame = caption_box.text_frame
            caption_frame.word_wrap = True
            caption_p = caption_frame.paragraphs[0]
            caption_p.text = alt_text
            caption_p.font.size = Pt(12)
            caption_p.font.italic = True
            caption_p.font.color.rgb = self._hex_to_rgb(colors.text)
            caption_p.alignment = PP_ALIGN.CENTER

    def _add_conclusion_slide(
        self,
        slide,
        content: Dict[str, Any],
        colors: ColorScheme
    ) -> None:
        """Add conclusion/takeaways slide"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        # Heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)
        )
        heading_frame = heading_box.text_frame
        heading_p = heading_frame.paragraphs[0]
        heading_p.text = content.get('heading', 'Key Takeaways')
        heading_p.font.size = Pt(40)
        heading_p.font.bold = True
        heading_p.font.color.rgb = self._hex_to_rgb(colors.primary)

        # Takeaways
        takeaways = content.get('takeaways', [])
        takeaway_y = 1.3
        for takeaway in takeaways[:5]:
            takeaway_box = slide.shapes.add_textbox(
                Inches(1), Inches(takeaway_y), Inches(8.5), Inches(0.5)
            )
            takeaway_frame = takeaway_box.text_frame
            takeaway_frame.word_wrap = True
            takeaway_p = takeaway_frame.paragraphs[0]
            takeaway_p.text = f"✓ {takeaway}"
            takeaway_p.font.size = Pt(16)
            takeaway_p.font.color.rgb = self._hex_to_rgb(colors.text)
            takeaway_y += 0.65

        # CTA
        cta_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(9), Inches(0.8)
        )
        cta_frame = cta_box.text_frame
        cta_frame.word_wrap = True
        cta_p = cta_frame.paragraphs[0]
        cta_p.text = content.get('cta', 'Thank you!')
        cta_p.font.size = Pt(14)
        cta_p.font.color.rgb = self._hex_to_rgb(colors.secondary)
        cta_p.alignment = PP_ALIGN.CENTER

    def _add_footer(
        self,
        slide,
        colors: ColorScheme,
        slide_num: int,
        total_slides: int
    ) -> None:
        """Add footer with slide number and author"""

        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # Slide number
        if self.config.include_slide_numbers:
            footer_box = slide.shapes.add_textbox(
                Inches(9), Inches(5.3), Inches(0.8), Inches(0.25)
            )
            footer_frame = footer_box.text_frame
            footer_p = footer_frame.paragraphs[0]
            footer_p.text = f"{slide_num}/{total_slides}"
            footer_p.font.size = Pt(8)
            footer_p.font.color.rgb = self._hex_to_rgb(colors.text)
            footer_p.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def _hex_to_rgb(hex_color: str):
        """Convert hex color to RGB object"""
        from pptx.dml.color import RGBColor

        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )


# Example usage
if __name__ == "__main__":
    config = PptxConfig(theme='technical')
    generator = PowerPointGenerator(config)

    # Example slides
    sample_slides = [
        {
            'type': 'title',
            'content': {
                'title': 'Hello PowerPoint',
                'author': 'Wallace Espindola',
                'date': '2025-02-02',
                'tags': ['python', 'presentations']
            }
        },
        {
            'type': 'content',
            'content': {
                'heading': 'Getting Started',
                'body': 'This is a sample presentation created from Python.',
                'bullet_points': [
                    'Easy to use',
                    'Automatic from articles',
                    'Beautiful themes'
                ]
            }
        }
    ]

    # This would need python-pptx installed
    # generator.generate_presentation(
    #     sample_slides,
    #     'presentations/sample/sample.pptx',
    #     {'title': 'Sample', 'author': 'Wallace'}
    # )
